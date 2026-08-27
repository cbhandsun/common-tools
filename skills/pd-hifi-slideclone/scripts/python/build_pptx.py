import argparse
import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


EMU_PER_PT = 12700


def emu(value):
    return int(round(float(value) * EMU_PER_PT))


def rgb(value, default="000000"):
    if not value:
        value = default
    value = value.strip().lstrip("#")
    if len(value) != 6:
        value = default
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def abs_path(base, value):
    if not value:
        return None
    return value if os.path.isabs(value) else os.path.abspath(os.path.join(base, value))


def add_textbox(slide, item):
    box = item["box"]
    shape = slide.shapes.add_textbox(emu(box["x"]), emu(box["y"]), emu(box["w"]), emu(box["h"]))
    apply_component_replacement_metadata(shape, item)
    if item.get("rotation") is not None:
        shape.rotation = float(item.get("rotation"))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = emu((item.get("style") or {}).get("marginLeftPt", 0))
    tf.margin_right = emu((item.get("style") or {}).get("marginRightPt", 0))
    tf.margin_top = emu((item.get("style") or {}).get("marginTopPt", 0))
    tf.margin_bottom = emu((item.get("style") or {}).get("marginBottomPt", 0))
    if (item.get("style") or {}).get("wrap") is False:
        tf.word_wrap = False
        set_text_frame_no_wrap(tf)
    if (item.get("font") or {}).get("valign") == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    align = (item.get("font") or {}).get("align")
    if align == "center":
        paragraph.alignment = PP_ALIGN.CENTER
    elif align == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    font = item.get("font") or {}
    effective_font = {**font, "opacity": effective_text_opacity(item)}
    apply_line_spacing(paragraph, font)
    lines = str(item.get("text", "")).splitlines() or [""]
    for index, line in enumerate(lines):
        current_paragraph = paragraph if index == 0 else tf.add_paragraph()
        current_paragraph.alignment = paragraph.alignment
        apply_line_spacing(current_paragraph, font)
        run = current_paragraph.add_run()
        run.text = line
        apply_run_font(run, effective_font)
    return shape


def effective_text_opacity(item):
    style = item.get("style") or {}
    if str(style.get("visibility", "")).lower() == "hidden":
        return 0
    if style.get("opacity") is not None:
        return style.get("opacity")
    return (item.get("font") or {}).get("opacity", 1)


def apply_run_font(run, font):
    run.font.size = Pt(font.get("sizePt", 14))
    run.font.bold = str(font.get("weight", "")).lower() == "bold"
    run.font.color.rgb = rgb(font.get("color"), "111111")
    apply_font_opacity(run, font.get("opacity", 1))
    if font.get("family"):
        run.font.name = font.get("family")
        set_run_typeface(run, font.get("family"))


def set_text_frame_no_wrap(text_frame):
    body_pr = text_frame._txBody.bodyPr
    body_pr.set("wrap", "none")
    for child in list(body_pr):
        if child.tag in {qn("a:spAutoFit"), qn("a:normAutofit"), qn("a:noAutofit")}:
            body_pr.remove(child)
    body_pr.append(OxmlElement("a:noAutofit"))


def apply_font_opacity(run, opacity):
    try:
        alpha_value = max(0, min(1, float(opacity)))
    except (TypeError, ValueError):
        alpha_value = 1
    if alpha_value >= 0.999:
        return
    r_pr = run._r.get_or_add_rPr()
    solid = r_pr.find(qn("a:solidFill"))
    if solid is None:
        solid = OxmlElement("a:solidFill")
        r_pr.append(solid)
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", "111111")
        solid.append(srgb)
    for child in list(srgb):
        if child.tag.endswith("}alpha"):
            srgb.remove(child)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(int(alpha_value * 100000)))
    srgb.append(alpha)


def apply_line_spacing(paragraph, font):
    line_height = font.get("lineHeightMultiple")
    if isinstance(line_height, (int, float)) and line_height > 0:
        paragraph.line_spacing = float(line_height)


def set_run_typeface(run, family):
    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag.endswith("}latin") or child.tag.endswith("}ea") or child.tag.endswith("}cs"):
            r_pr.remove(child)
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = OxmlElement(tag)
        node.set("typeface", family)
        r_pr.append(node)


def shape_type(kind):
    kind = (kind or "rect").lower()
    if kind in ("ellipse", "oval", "circle"):
        return MSO_SHAPE.OVAL
    if kind in ("roundrect", "rounded-rect", "roundedrectangle", "phone", "mobile", "device-phone"):
        return MSO_SHAPE.ROUNDED_RECTANGLE
    if kind == "triangle":
        return MSO_SHAPE.ISOSCELES_TRIANGLE
    if kind in ("right-triangle", "righttriangle"):
        return MSO_SHAPE.RIGHT_TRIANGLE
    if kind == "diamond":
        return MSO_SHAPE.DIAMOND
    if kind == "hexagon":
        return MSO_SHAPE.HEXAGON
    if kind == "chevron":
        return MSO_SHAPE.CHEVRON
    if kind == "parallelogram":
        return MSO_SHAPE.PARALLELOGRAM
    if kind == "cylinder":
        return MSO_SHAPE.CAN
    if kind == "cloud":
        return MSO_SHAPE.CLOUD
    if kind in ("document", "flowchart-document", "flowchartdocument"):
        return MSO_SHAPE.FLOWCHART_DOCUMENT
    if kind in ("screen", "device-screen", "monitor"):
        return MSO_SHAPE.FRAME
    if kind in ("funnel", "filter-funnel"):
        return MSO_SHAPE.FUNNEL
    if kind in ("donut", "ring"):
        return MSO_SHAPE.DONUT
    if kind == "arc":
        return MSO_SHAPE.ARC
    if kind in ("blockarc", "block-arc"):
        return MSO_SHAPE.BLOCK_ARC
    if kind in ("circulararrow", "circular-arrow", "cycle-arrow"):
        return MSO_SHAPE.CIRCULAR_ARROW
    if kind in ("bentarrow", "bent-arrow"):
        return MSO_SHAPE.BENT_ARROW
    if kind in ("leftarrow", "left-arrow"):
        return MSO_SHAPE.LEFT_ARROW
    if kind in ("rightarrow", "right-arrow"):
        return MSO_SHAPE.RIGHT_ARROW
    if kind in ("uparrow", "up-arrow"):
        return MSO_SHAPE.UP_ARROW
    if kind in ("downarrow", "down-arrow"):
        return MSO_SHAPE.DOWN_ARROW
    if kind in ("leftrightarrow", "left-right-arrow"):
        return MSO_SHAPE.LEFT_RIGHT_ARROW
    if kind in ("updownarrow", "up-down-arrow"):
        return MSO_SHAPE.UP_DOWN_ARROW
    if kind in ("curvedleftarrow", "curved-left-arrow"):
        return MSO_SHAPE.CURVED_LEFT_ARROW
    if kind in ("curvedrightarrow", "curved-right-arrow"):
        return MSO_SHAPE.CURVED_RIGHT_ARROW
    if kind in ("uturnarrow", "u-turn-arrow"):
        return MSO_SHAPE.UTURN_ARROW
    if kind == "line":
        return None
    return MSO_SHAPE.RECTANGLE


def add_shape(slide, item, shape_index=None):
    box = item["box"]
    style = item.get("style") or {}
    if (item.get("type") or "").lower() == "line":
        connector_type = connector_type_for(style.get("connectorType"))
        start_x, start_y = anchor_point(style.get("startAnchor"), shape_index, box["x"], box["y"])
        end_x, end_y = anchor_point(style.get("endAnchor"), shape_index, box["x"] + box["w"], box["y"] + box["h"])
        shape = slide.shapes.add_connector(connector_type, emu(start_x), emu(start_y), emu(end_x), emu(end_y))
    elif (item.get("type") or "").lower() in ("freeform", "polyline"):
        points = item.get("points") or []
        if len(points) < 3:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(box["x"]), emu(box["y"]), emu(box["w"]), emu(box["h"]))
        else:
            local_points = [(emu(point["x"] * box["w"]), emu(point["y"] * box["h"])) for point in points]
            builder = slide.shapes.build_freeform(local_points[0][0], local_points[0][1])
            builder.add_line_segments(local_points[1:], close=should_close_freeform(item, style))
            shape = builder.convert_to_shape(emu(box["x"]), emu(box["y"]))
    else:
        shape = slide.shapes.add_shape(shape_type(item.get("type")), emu(box["x"]), emu(box["y"]), emu(box["w"]), emu(box["h"]))
        if (item.get("type") or "").lower() in ("roundrect", "rounded-rect", "roundedrectangle", "phone", "mobile", "device-phone") and style.get("radiusRatio") is not None:
            try:
                shape.adjustments[0] = float(style.get("radiusRatio"))
            except Exception:
                pass
        apply_shape_adjustments(shape, style)
        rotation = item.get("rotation") if item.get("rotation") is not None else style.get("rotation")
        if rotation is not None:
            shape.rotation = float(rotation)
    if (item.get("type") or "").lower() != "line":
        apply_shape_fill(shape, style)
    stroke = style.get("stroke")
    stroke_width = style.get("strokeWidthPt")
    if stroke == "none" or (stroke_width is not None and float(stroke_width) <= 0):
        shape.line.fill.background()
    elif stroke:
        shape.line.color.rgb = rgb(stroke, "000000")
        apply_line_opacity(shape, style)
    if stroke_width is not None and float(stroke_width) > 0:
        shape.line.width = Pt(stroke_width)
    if style.get("dash"):
        add_line_dash(shape, style.get("dash"))
    if style.get("endArrow"):
        add_line_end(shape, style.get("endArrow"), "tailEnd")
    if style.get("startArrow"):
        add_line_end(shape, style.get("startArrow"), "headEnd")
    if style.get("shadow"):
        add_outer_shadow(shape, style.get("shadow"))
    else:
        clear_shape_effects(shape)
    apply_component_replacement_metadata(shape, item)
    return shape


def should_close_freeform(item, style):
    if (item.get("type") or "").lower() == "polyline":
        return False
    if style.get("closePath") is not None:
        return bool(style.get("closePath"))
    return True


def apply_shape_adjustments(shape, style):
    adjustments = style.get("adjustments")
    if not isinstance(adjustments, list):
        return
    for index, value in enumerate(adjustments):
        try:
            if index < len(shape.adjustments) and value is not None:
                shape.adjustments[index] = float(value)
        except Exception:
            pass


def apply_shape_fill(shape, style):
    gradient = style.get("gradient")
    if is_valid_gradient(gradient):
        add_gradient_fill(shape, gradient)
        return
    fill = style.get("fill")
    if fill and str(fill).lower() != "none":
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill, "FFFFFF")
        apply_solid_fill_opacity(shape, style)
    else:
        shape.fill.background()


def opacity_value(style, specific_key):
    if style.get(specific_key) is not None:
        return clamp_float(style.get(specific_key), 0, 1)
    if style.get("opacity") is not None:
        return clamp_float(style.get("opacity"), 0, 1)
    return None


def apply_solid_fill_opacity(shape, style):
    alpha_value = opacity_value(style, "fillOpacity")
    if alpha_value is None or alpha_value >= 0.999:
        return
    sp_pr = shape.element.spPr
    solid = sp_pr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    set_alpha_node(srgb, alpha_value)


def apply_line_opacity(shape, style):
    alpha_value = opacity_value(style, "strokeOpacity")
    if alpha_value is None or alpha_value >= 0.999:
        return
    line = shape.element.spPr.find(qn("a:ln"))
    if line is None:
        return
    solid = line.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    set_alpha_node(srgb, alpha_value)


def set_alpha_node(srgb, alpha_value):
    for child in list(srgb):
        if child.tag.endswith("}alpha"):
            srgb.remove(child)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(int(clamp_float(alpha_value, 0, 1) * 100000)))
    srgb.append(alpha)


def is_valid_gradient(gradient):
    if not isinstance(gradient, dict):
        return False
    stops = gradient.get("stops")
    return isinstance(stops, list) and len(stops) >= 2


def add_gradient_fill(shape, gradient):
    sp_pr = shape.element.spPr
    remove_fill_nodes(sp_pr)
    grad_fill = OxmlElement("a:gradFill")
    grad_fill.set("flip", "none")
    grad_fill.set("rotWithShape", "1")
    gs_lst = OxmlElement("a:gsLst")
    stops = normalize_gradient_stops(gradient.get("stops") or [])
    for stop in stops:
        gs = OxmlElement("a:gs")
        gs.set("pos", str(stop["position"]))
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", normalize_hex(stop["color"], "FFFFFF"))
        alpha_value = stop.get("alpha")
        if alpha_value is not None:
            alpha = OxmlElement("a:alpha")
            alpha.set("val", str(int(clamp_float(alpha_value, 0, 1) * 100000)))
            srgb.append(alpha)
        gs.append(srgb)
        gs_lst.append(gs)
    grad_fill.append(gs_lst)
    lin = OxmlElement("a:lin")
    lin.set("ang", str(int(float(gradient.get("angleDeg", 0) or 0) * 60000)))
    lin.set("scaled", "1")
    grad_fill.append(lin)
    insert_fill_node(sp_pr, grad_fill)


def normalize_gradient_stops(stops):
    normalized = []
    for index, stop in enumerate(stops):
        if not isinstance(stop, dict):
            continue
        color = stop.get("color")
        if not color:
            continue
        raw_position = stop.get("position")
        if raw_position is None:
            raw_position = index / max(1, len(stops) - 1)
        position = int(round(clamp_float(raw_position, 0, 1) * 100000))
        normalized.append({
            "color": color,
            "position": position,
            "alpha": stop.get("alpha")
        })
    if len(normalized) < 2:
        normalized = [
            {"color": "#FFFFFF", "position": 0, "alpha": None},
            {"color": "#FFFFFF", "position": 100000, "alpha": None}
        ]
    normalized.sort(key=lambda item: item["position"])
    normalized[0]["position"] = 0
    normalized[-1]["position"] = 100000
    return normalized


def remove_fill_nodes(sp_pr):
    fill_tags = {qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"), qn("a:blipFill"), qn("a:pattFill"), qn("a:grpFill")}
    for child in list(sp_pr):
        if child.tag in fill_tags:
            sp_pr.remove(child)


def insert_fill_node(sp_pr, fill_node):
    for index, child in enumerate(list(sp_pr)):
        if child.tag == qn("a:ln") or child.tag.endswith("}effectLst"):
            sp_pr.insert(index, fill_node)
            return
    sp_pr.append(fill_node)


def clamp_float(value, minimum, maximum):
    try:
        number = float(value)
    except (TypeError, ValueError):
        number = minimum
    return max(minimum, min(maximum, number))


def clear_shape_effects(shape):
    sp_pr = shape.element.spPr
    for child in list(sp_pr):
        if child.tag.endswith("}effectLst"):
            sp_pr.remove(child)
    sp_pr.append(OxmlElement("a:effectLst"))


def add_outer_shadow(shape, shadow):
    sp_pr = shape.element.spPr
    for child in list(sp_pr):
        if child.tag.endswith("}effectLst"):
            sp_pr.remove(child)
    effect_lst = OxmlElement("a:effectLst")
    outer = OxmlElement("a:outerShdw")
    outer.set("blurRad", str(emu(shadow.get("blurPt", 4))))
    outer.set("dist", str(emu(shadow.get("distancePt", 1.5))))
    outer.set("dir", str(int(float(shadow.get("angleDeg", 45)) * 60000)))
    outer.set("algn", "ctr")
    outer.set("rotWithShape", "0")
    color = OxmlElement("a:srgbClr")
    color_value = (shadow.get("color") or "#000000").strip().lstrip("#")
    if len(color_value) != 6:
        color_value = "000000"
    color.set("val", color_value)
    alpha = OxmlElement("a:alpha")
    alpha_value = max(0, min(1, float(shadow.get("alpha", 0.18))))
    alpha.set("val", str(int(alpha_value * 100000)))
    color.append(alpha)
    outer.append(color)
    effect_lst.append(outer)
    sp_pr.append(effect_lst)


def connector_type_for(value):
    value = (value or "straight").lower()
    if value == "elbow":
        return MSO_CONNECTOR.ELBOW
    if value == "curve":
        return MSO_CONNECTOR.CURVE
    return MSO_CONNECTOR.STRAIGHT


def add_line_end(shape, arrow_type, tag_name):
    sp_pr = shape.element.spPr
    ln = sp_pr.ln
    if ln is None:
        ln = OxmlElement("a:ln")
        sp_pr.append(ln)
    for child in list(ln):
        if child.tag.endswith("}" + tag_name):
            ln.remove(child)
    end = OxmlElement(f"a:{tag_name}")
    end.set("type", arrow_type if isinstance(arrow_type, str) else "triangle")
    ln.append(end)


def add_line_dash(shape, dash_type):
    sp_pr = shape.element.spPr
    ln = sp_pr.ln
    if ln is None:
        ln = OxmlElement("a:ln")
        sp_pr.append(ln)
    for child in list(ln):
        if child.tag.endswith("}prstDash"):
            ln.remove(child)
    dash = OxmlElement("a:prstDash")
    dash_value = str(dash_type or "dash").lower()
    dash.set("val", "dash" if dash_value in ("dash", "dashed") else dash_value)
    ln.append(dash)


def add_picture(slide, item, base):
    if item.get("type") == "source-reference" and (item.get("style") or {}).get("opacity", 1) < 0.99:
        return None
    source = item.get("source") or {}
    style = item.get("style") or {}
    asset = item.get("assetPath") or (item.get("style") or {}).get("assetPath") or source.get("pageImage")
    asset = abs_path(base, asset)
    if not asset or not os.path.exists(asset):
        return None
    box = item["box"]
    picture = slide.shapes.add_picture(asset, emu(box["x"]), emu(box["y"]), emu(box["w"]), emu(box["h"]))
    apply_picture_crop(picture, style.get("crop"))
    apply_component_replacement_metadata(picture, item)
    return picture


def apply_picture_crop(picture, crop):
    if not isinstance(crop, dict):
        return
    crop_map = {
        "left": "crop_left",
        "top": "crop_top",
        "right": "crop_right",
        "bottom": "crop_bottom"
    }
    for key, attr in crop_map.items():
        if key not in crop:
            continue
        try:
            value = max(0.0, min(1.0, float(crop.get(key))))
            setattr(picture, attr, value)
        except Exception:
            pass


def add_table(slide, item):
    rows = item.get("rows") or []
    if not rows:
        return None
    col_count = max(len(row) for row in rows)
    box = item["box"]
    shape = slide.shapes.add_table(len(rows), col_count, emu(box["x"]), emu(box["y"]), emu(box["w"]), emu(box["h"]))
    apply_component_replacement_metadata(shape, item)
    table = shape.table
    style = item.get("style") or {}
    configure_table_style(table)
    column_widths = normalized_table_dimensions(style.get("columnWidthsPt"), col_count, box["w"])
    row_heights = normalized_table_dimensions(style.get("rowHeightsPt"), len(rows), box["h"])
    for column_index, column in enumerate(table.columns):
        column.width = emu(column_widths[column_index])
    for row_index, row in enumerate(table.rows):
        row.height = emu(row_heights[row_index])
    for row_index, row in enumerate(rows):
        for col_index in range(col_count):
            cell = table.cell(row_index, col_index)
            cell.text = "" if uses_overlay_text(style) else (row[col_index] if col_index < len(row) else "")
            style_table_cell(cell, style, is_header=(row_index == 0), cell_style=table_cell_style(style, row_index, col_index))
    if uses_overlay_grid(style):
        add_table_grid_overlays(slide, item, col_count)
    if uses_overlay_text(style):
        add_table_text_overlays(slide, item, col_count)
    return shape


def normalized_table_dimensions(values, count, total):
    if not isinstance(count, int) or count <= 0:
        return []
    if not isinstance(values, list) or len(values) != count:
        return [float(total) / count] * count
    normalized = []
    for value in values:
        if isinstance(value, bool):
            return [float(total) / count] * count
        try:
            numeric = float(value)
        except (TypeError, ValueError):
            return [float(total) / count] * count
        if numeric <= 0 or numeric != numeric or numeric > 1000000:
            return [float(total) / count] * count
        normalized.append(numeric)
    dimension_sum = sum(normalized)
    if dimension_sum <= 0:
        return [float(total) / count] * count
    scale = float(total) / dimension_sum
    return [value * scale for value in normalized]


def table_cell_style(style, row_index, col_index):
    matrix = style.get("cellStyles")
    if not isinstance(matrix, list) or row_index < 0 or row_index >= len(matrix):
        return {}
    row = matrix[row_index]
    if not isinstance(row, list) or col_index < 0 or col_index >= len(row):
        return {}
    return row[col_index] if isinstance(row[col_index], dict) else {}


def anchor_point(anchor, shape_index, fallback_x, fallback_y):
    if not isinstance(anchor, dict) or not shape_index:
        return fallback_x, fallback_y
    target_id = anchor.get("elementId") or anchor.get("id")
    target = shape_index.get(target_id)
    box = (target or {}).get("box") or {}
    if not all(is_number(box.get(key)) for key in ("x", "y", "w", "h")):
        return fallback_x, fallback_y
    position = clamp_number(anchor.get("position"), 0.5, 0, 1)
    side = str(anchor.get("side") or "center").lower()
    x = float(box["x"])
    y = float(box["y"])
    w = float(box["w"])
    h = float(box["h"])
    if side == "left":
        point = (x, y + h * position)
    elif side == "right":
        point = (x + w, y + h * position)
    elif side == "top":
        point = (x + w * position, y)
    elif side == "bottom":
        point = (x + w * position, y + h)
    else:
        point = (x + w / 2, y + h / 2)
    return (
        point[0] + float(anchor.get("dxPt") or 0),
        point[1] + float(anchor.get("dyPt") or 0)
    )


def is_number(value):
    return isinstance(value, (int, float)) and not isinstance(value, bool)


def clamp_number(value, default, minimum, maximum):
    if not is_number(value):
        return default
    return max(minimum, min(maximum, float(value)))


def uses_overlay_text(style):
    return str(style.get("textMode", "")).lower() in ("overlay-textboxes", "textboxes")


def uses_overlay_grid(style):
    return str(style.get("gridMode", "")).lower() in ("overlay-lines", "lines")


def add_table_grid_overlays(slide, item, col_count):
    rows = item.get("rows") or []
    if not rows:
        return
    style = item.get("style") or {}
    box = item["box"]
    stroke = style.get("stroke", "#D0D7DE")
    stroke_width = style.get("strokeWidthPt", 0.35)
    row_height = box["h"] / max(1, len(rows))
    col_width = box["w"] / col_count
    for row_index in range(len(rows) + 1):
        y = box["y"] + row_height * row_index
        add_shape(slide, {
            "id": "%s-grid-h%s" % (item.get("id", "table"), row_index),
            "type": "line",
            "box": { "x": box["x"], "y": y, "w": box["w"], "h": 0 },
            "style": { "stroke": stroke, "strokeWidthPt": stroke_width }
        })
    for col_index in range(col_count + 1):
        x = box["x"] + col_width * col_index
        add_shape(slide, {
            "id": "%s-grid-v%s" % (item.get("id", "table"), col_index),
            "type": "line",
            "box": { "x": x, "y": box["y"], "w": 0, "h": box["h"] },
            "style": { "stroke": stroke, "strokeWidthPt": stroke_width }
        })


def add_table_text_overlays(slide, item, col_count):
    rows = item.get("rows") or []
    style = item.get("style") or {}
    box = item["box"]
    col_width = box["w"] / col_count
    row_height = box["h"] / max(1, len(rows))
    padding_left = style.get("textBoxPaddingLeftPt", style.get("paddingLeftPt", 5))
    padding_right = style.get("textBoxPaddingRightPt", style.get("paddingRightPt", 5))
    padding_top = style.get("textBoxPaddingTopPt", 0)
    padding_bottom = style.get("textBoxPaddingBottomPt", 0)
    for row_index, row in enumerate(rows):
        is_header = row_index == 0
        for col_index in range(col_count):
            text_value = row[col_index] if col_index < len(row) else ""
            text_box = {
                "id": "%s-r%s-c%s-text" % (item.get("id", "table"), row_index, col_index),
                "text": text_value,
                "box": {
                    "x": box["x"] + col_width * col_index + padding_left,
                    "y": box["y"] + row_height * row_index + padding_top,
                    "w": max(1, col_width - padding_left - padding_right),
                    "h": max(1, row_height - padding_top - padding_bottom)
                },
                "font": {
                    "family": style.get("fontFamily"),
                    "sizePt": style.get("headerFontSizePt") if is_header and style.get("headerFontSizePt") else style.get("fontSizePt", 14),
                    "weight": style.get("headerWeight") if is_header and style.get("headerWeight") else style.get("fontWeight", "regular"),
                    "color": style.get("headerTextColor") if is_header and style.get("headerTextColor") else style.get("textColor", "#111111"),
                    "align": style.get("textAlign", "left"),
                    "valign": style.get("textValign", "middle")
                },
                "style": {
                    "marginLeftPt": 0,
                    "marginRightPt": 0,
                    "marginTopPt": 0,
                    "marginBottomPt": 0
                }
            }
            if is_header and not text_box["font"].get("weight"):
                text_box["font"]["weight"] = "bold"
            add_textbox(slide, text_box)


def configure_table_style(table):
    table._tbl.tblPr.set("firstRow", "0")
    table._tbl.tblPr.set("bandRow", "0")
    style_id = table._tbl.tblPr.find(qn("a:tableStyleId"))
    if style_id is not None:
      style_id.text = "{00000000-0000-0000-0000-000000000000}"


def style_table_cell(cell, style, is_header=False, cell_style=None):
    cell_style = cell_style if isinstance(cell_style, dict) else {}
    fill_color = cell_style.get("fill") or (style.get("headerFill") if is_header and style.get("headerFill") else style.get("fill", "#FFFFFF"))
    text_color = cell_style.get("textColor") or (style.get("headerTextColor") if is_header and style.get("headerTextColor") else style.get("textColor", "#111111"))
    border_color = cell_style.get("stroke") or style.get("stroke", "#D0D7DE")
    font_family = cell_style.get("fontFamily") or style.get("fontFamily")
    font_size = cell_style.get("fontSizePt") or (style.get("headerFontSizePt") if is_header and style.get("headerFontSizePt") else style.get("fontSizePt", 14))
    font_weight = cell_style.get("fontWeight") or (style.get("headerWeight") if is_header and style.get("headerWeight") else style.get("fontWeight"))
    if not font_weight and is_header:
        font_weight = "bold"

    if str(fill_color).lower() == "none":
        cell.fill.background()
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill_color, "FFFFFF")
    cell.margin_left = emu(cell_style.get("paddingLeftPt", style.get("paddingLeftPt", 8)))
    cell.margin_right = emu(cell_style.get("paddingRightPt", style.get("paddingRightPt", 8)))
    cell.margin_top = emu(cell_style.get("paddingTopPt", style.get("paddingTopPt", 4)))
    cell.margin_bottom = emu(cell_style.get("paddingBottomPt", style.get("paddingBottomPt", 4)))

    text_frame = cell.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE if cell_style.get("textValign", style.get("textValign", "middle")) == "middle" else MSO_ANCHOR.TOP
    for paragraph in text_frame.paragraphs:
        alignment = cell_style.get("textAlign", style.get("textAlign", "left"))
        paragraph.alignment = PP_ALIGN.CENTER if alignment == "center" else PP_ALIGN.RIGHT if alignment == "right" else PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = str(font_weight or "").lower() == "bold"
            run.font.color.rgb = rgb(text_color, "111111")
            if font_family:
                run.font.name = font_family
                set_run_typeface(run, font_family)

    border_color = "none" if uses_overlay_grid(style) else border_color
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        set_cell_border(cell, edge, border_color, style.get("strokeWidthPt", 0.6))


def set_cell_border(cell, edge_tag, color_value, width_pt):
    tc_pr = cell._tc.get_or_add_tcPr()
    edge = tc_pr.find(qn(edge_tag))
    if edge is None:
        edge = OxmlElement(edge_tag)
        tc_pr.append(edge)
    edge.set("w", str(emu(width_pt)))
    edge.set("cap", "flat")
    edge.set("cmpd", "sng")
    edge.set("algn", "ctr")
    clear_children(edge)
    if str(color_value).lower() == "none" or float(width_pt or 0) <= 0:
        edge.append(OxmlElement("a:noFill"))
        return
    solid_fill = OxmlElement("a:solidFill")
    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", normalize_hex(color_value, "D0D7DE"))
    solid_fill.append(srgb)
    edge.append(solid_fill)
    edge.append(OxmlElement("a:prstDash"))
    edge[-1].set("val", "solid")
    edge.append(OxmlElement("a:round"))
    edge.append(OxmlElement("a:headEnd"))
    edge[-1].set("type", "none")
    edge.append(OxmlElement("a:tailEnd"))
    edge[-1].set("type", "none")


def clear_children(node):
    for child in list(node):
        node.remove(child)


def normalize_hex(value, default="000000"):
    if not value:
        value = default
    value = str(value).strip().lstrip("#")
    if len(value) != 6:
        value = default
    return value.upper()


def apply_component_replacement_metadata(shape, item):
    description = build_component_replacement_description(item.get("source") or item.get("style") or {})
    if not description:
        return
    try:
        if hasattr(shape, "name") and item.get("id"):
            shape.name = str(item.get("id"))[:255]
    except Exception:
        pass
    c_nv_pr = find_c_nv_pr(shape)
    if c_nv_pr is not None:
        c_nv_pr.set("descr", description)


def find_c_nv_pr(shape):
    element = getattr(shape, "element", None)
    if element is None:
        return None
    for node in element.iter():
        if str(node.tag).endswith("}cNvPr"):
            return node
    return None


def build_component_replacement_description(metadata):
    if not isinstance(metadata, dict):
        return None
    plan = metadata.get("componentReplacementPlan")
    if not isinstance(plan, dict):
        return None
    source_provider = plan.get("sourceProvider")
    component_id = plan.get("componentId") or metadata.get("componentReplacementCandidateId")
    if not source_provider or not component_id:
        return None
    parts = [
        "slideclone:componentReplacementPlan",
        "provider=%s" % sanitize_metadata_value(source_provider, 48),
        "kind=%s" % sanitize_metadata_value(plan.get("componentKind") or "component", 32),
        "id=%s" % sanitize_metadata_value(component_id, 96)
    ]
    layer_key = plan.get("layerKey") or metadata.get("componentReplacementLayerKey")
    tier = plan.get("suitabilityTier") or metadata.get("componentReplacementSuitabilityTier")
    score = plan.get("suitabilityScore", metadata.get("componentReplacementSuitabilityScore"))
    if layer_key:
        parts.append("layer=%s" % sanitize_metadata_value(layer_key, 48))
    if tier:
        parts.append("tier=%s" % sanitize_metadata_value(tier, 32))
    if isinstance(score, (int, float)) and not isinstance(score, bool):
        parts.append("score=%s" % ("%0.2f" % float(score)).rstrip("0").rstrip("."))
    return " ".join(parts)


def sanitize_metadata_value(value, max_length):
    text = str(value)
    chars = []
    for char in text:
        if len(chars) >= max_length:
            break
        chars.append("_" if char.isspace() or ord(char) < 32 else char)
    return "".join(chars).strip("_")


def qn(name):
    prefix, tagroot = name.split(":")
    namespaces = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main"
    }
    return "{%s}%s" % (namespaces.get(prefix, namespaces["a"]), tagroot)


def build(ir_file, out_file):
    with open(ir_file, "r", encoding="utf-8") as handle:
        deck = json.load(handle)
    base = os.path.dirname(os.path.abspath(ir_file))

    prs = Presentation()
    prs.slide_width = emu(deck["slideSize"]["widthPt"])
    prs.slide_height = emu(deck["slideSize"]["heightPt"])
    blank_layout = prs.slide_layouts[6]

    for page in sorted(deck.get("pages", []), key=lambda item: item.get("pageIndex", 0)):
        slide = prs.slides.add_slide(blank_layout)
        background = page.get("background") or {}
        if background.get("fill"):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(background.get("fill"), "FFFFFF")
            shape.line.fill.background()
        shape_index = page_shape_index(page)
        table_overlay_images = [
            item for item in page.get("images", [])
            if item.get("source", {}).get("tableOverlay") is True
        ]
        leading_images = [
            item for item in page.get("images", [])
            if not item.get("drawAfterShapes") and item not in table_overlay_images
        ]
        trailing_images = [
            item for item in page.get("images", [])
            if item.get("drawAfterShapes") and item not in table_overlay_images
        ]
        for item in leading_images:
            add_picture(slide, item, base)
        for item in page.get("shapes", []):
            add_shape(slide, item, shape_index)
        for item in trailing_images:
            add_picture(slide, item, base)
        for item in page.get("tables", []):
            add_table(slide, item)
        for item in table_overlay_images:
            add_picture(slide, item, base)
        for item in page.get("textBoxes", []):
            add_textbox(slide, item)

    os.makedirs(os.path.dirname(os.path.abspath(out_file)), exist_ok=True)
    prs.save(out_file)


def page_shape_index(page):
    items = []
    for collection_name in ("shapes", "images", "tables", "textBoxes"):
        items.extend(page.get(collection_name, []))
    return {
        item.get("id"): item
        for item in items
        if isinstance(item, dict) and item.get("id") and isinstance(item.get("box"), dict)
    }


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--ir", required=True)
    parser.add_argument("--out", required=True)
    args = parser.parse_args()
    build(args.ir, args.out)


if __name__ == "__main__":
    main()
